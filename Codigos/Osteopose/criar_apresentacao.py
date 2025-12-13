from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# 1. Criar a Apresentação
prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)

# Função auxiliar para definir cores Hex
def set_color(shape, hex_color):
    rgb = RGBColor.from_string(hex_color)
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb
    shape.line.fill.background() # Sem borda

# Função para texto
def add_text(slide, text, left, top, width, height, font_size, is_bold=False, color='000000', align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = is_bold
    p.font.name = 'Arial'
    p.font.color.rgb = RGBColor.from_string(color)
    p.alignment = align
    return txBox

# ==============================================================================
# SLIDE 1: Abertura Técnica (Fundo Azul)
# ==============================================================================
slide1 = prs.slides.add_slide(prs.slide_layouts[6]) # Layout vazio
background = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(16), Inches(9))
set_color(background, '0a192f') # Azul Marinho Profundo

# Títulos
add_text(slide1, "Mudanças de Paradigma", Inches(0), Inches(2.5), Inches(16), Inches(1), 60, True, 'FFFFFF', PP_ALIGN.CENTER)
add_text(slide1, "no Tratamento da Osteoporose", Inches(0), Inches(3.5), Inches(16), Inches(1), 60, True, 'FFFFFF', PP_ALIGN.CENTER)
add_text(slide1, "Do T-score ao Treat-to-Target", Inches(0), Inches(5), Inches(16), Inches(1), 32, False, 'D4AF37', PP_ALIGN.CENTER)




# ==============================================================================
# SLIDE 2: O Arco (Fundo Claro - Simétrico)
# ==============================================================================
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
bg2 = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(16), Inches(9))
set_color(bg2, 'F9F9F9') # Off-white

# Título do Slide
add_text(slide2, "Objetivos Educacionais", Inches(1), Inches(0.5), Inches(10), Inches(1), 40, True, '0a192f', PP_ALIGN.LEFT)

# Desenhando o Arco (Simulando com uma forma de Arco)
arc = slide2.shapes.add_shape(MSO_SHAPE.ARC, Inches(2), Inches(4), Inches(12), Inches(5))
arc.line.color.rgb = RGBColor.from_string('D4AF37') # Dourado
arc.line.width = Pt(4) # Linha mais fina e elegante
arc.adjustments[0] = 180 # Começa na esquerda
arc.adjustments[1] = 0   # Termina na direita

# Dados dos Pontos
pontos = [
    {"num": "1", "txt": "Modelos de Risco & Limitações", "x": 2.5, "y": 7.5},
    {"num": "2", "txt": "A Nova Era dos Limiares?", "x": 4.5, "y": 6.0},
    {"num": "3", "txt": "Fratura Iminente", "x": 8.0, "y": 5.0}, # O Topo (Centralizado)
    {"num": "4", "txt": "Treat to target", "x": 11.5, "y": 6.0},
    {"num": "5", "txt": "Anabolic First", "x": 13.5, "y": 7.5},
]

for p in pontos:
    # Círculo Azul
    circle = slide2.shapes.add_shape(MSO_SHAPE.OVAL, Inches(p["x"]-0.35), Inches(p["y"]-0.35), Inches(0.7), Inches(0.7))
    set_color(circle, '0a192f')
    
    # Número Branco
    add_text(slide2, p["num"], Inches(p["x"]-0.35), Inches(p["y"]-0.25), Inches(0.7), Inches(0.7), 24, True, 'FFFFFF', PP_ALIGN.CENTER)
    
    # Texto Descritivo (Ajuste de posição para simetria)
    if p["num"] == "3": # Central
        add_text(slide2, p["txt"], Inches(p["x"]-2), Inches(p["y"]-1.0), Inches(4), Inches(0.5), 18, True, '333333', PP_ALIGN.CENTER)
    elif p["num"] in ["1", "2"]: # Lado Esquerdo -> Texto à Direita
        add_text(slide2, p["txt"], Inches(p["x"]+0.5), Inches(p["y"]-0.25), Inches(4), Inches(0.5), 18, False, '333333', PP_ALIGN.LEFT)
    else: # Lado Direito -> Texto à Esquerda (para balancear) ou Direita (se preferir fluxo)
        # Vamos manter à direita para leitura ocidental, mas alinhado
        add_text(slide2, p["txt"], Inches(p["x"]+0.5), Inches(p["y"]-0.25), Inches(4), Inches(0.5), 18, False, '333333', PP_ALIGN.LEFT)




# ==============================================================================
# SLIDE 3: Rilke (Círculos Crescentes)
# ==============================================================================
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
# Nota para você inserir a imagem
note = add_text(slide3, "[INSIRA AQUI A IMAGEM AZUL/DOURADA DOS CÍRCULOS]", Inches(2), Inches(4), Inches(12), Inches(1), 24, True, '000000', PP_ALIGN.CENTER)




# ==============================================================================
# SLIDE 4: Rilke (Dragões)
# ==============================================================================
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
# Nota para inserir imagem
add_text(slide4, "[INSIRA AQUI A IMAGEM DA PAISAGEM DOURADA]", Inches(2), Inches(2), Inches(12), Inches(1), 24, True, '000000', PP_ALIGN.CENTER)

# Texto da Poesia (Já diagramado para ficar no centro)
quote = "Talvez todos os dragões de nossa vida sejam princesas\nque esperam apenas ver-nos\nbelos e corajosos."

add_text(slide4, quote, Inches(3), Inches(4), Inches(10), Inches(3), 28, True, 'FFFFFF', PP_ALIGN.CENTER)
add_text(slide4, "— Rainer Maria Rilke", Inches(3), Inches(6.5), Inches(10), Inches(1), 18, False, 'D4AF37', PP_ALIGN.CENTER)




# Salvar
prs.save('Apresentacao_Curadoria.pptx')
print("Arquivo 'Apresentacao_Curadoria.pptx' gerado com sucesso!")


