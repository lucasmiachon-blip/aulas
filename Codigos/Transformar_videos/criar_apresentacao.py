from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# 1. Definindo a Paleta Turner (RGB)
# O código não "vê" cores, ele vê números.
OFF_WHITE = RGBColor(246, 242, 236)  # Fundo Claro
CARVAO    = RGBColor(42, 38, 36)     # Fundo Escuro
TERRACOTA = RGBColor(142, 59, 47)    # Destaque Fogo
OURO      = RGBColor(184, 137, 74)   # Destaque Montaigne
AZUL      = RGBColor(31, 58, 68)     # Destaque Dados

# Criando a apresentação vazia
prs = Presentation()

# --- SLIDE 1: ABERTURA (Turner) ---
slide = prs.slides.add_slide(prs.slide_layouts[6]) # Layout vazio
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = CARVAO # Fundo escuro (simulando a pintura)

# Título
box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(2))
tf = box.text_frame
p = tf.paragraphs[0]
p.text = "A CORAGEM NA INCERTEZA"
p.font.size = Pt(54)
p.font.name = 'Garamond'
p.font.color.rgb = OFF_WHITE
p.font.bold = True

# Subtítulo
box_sub = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(1))
p_sub = box_sub.text_frame.paragraphs[0]
p_sub.text = "GRADE aplicado à Diretriz Brasileira de Dislipidemias 2025"
p_sub.font.size = Pt(24)
p_sub.font.name = 'Arial'
p_sub.font.color.rgb = OURO

# --- SLIDE 2: MONTAIGNE (Silêncio) ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = CARVAO # Fundo Carvão

# Citação
box = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(6), Inches(2))
tf = box.text_frame
p = tf.paragraphs[0]
p.text = "Que sais-je?"
p.font.size = Pt(40)
p.font.name = 'Garamond'
p.font.italic = True
p.font.color.rgb = OURO
p.alignment = 2 # Centralizado

# Autor
p2 = tf.add_paragraph()
p2.text = "— Montaigne"
p2.font.size = Pt(18)
p2.font.color.rgb = OFF_WHITE
p2.alignment = 2

# --- SLIDE 3: DRUG HOLIDAY (Conteúdo) ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = OFF_WHITE # Fundo Claro

# Título do Slide
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
p = title_box.text_frame.paragraphs[0]
p.text = "Drug Holiday: Critérios de Segurança"
p.font.size = Pt(32)
p.font.color.rgb = AZUL
p.font.bold = True

# Coluna Vermelha (Quem NÃO pode)
shape = slide.shapes.add_shape(1, Inches(5.5), Inches(2), Inches(4), Inches(4)) # 1 = Retângulo
shape.fill.solid()
shape.fill.fore_color.rgb = TERRACOTA
shape.text_frame.text = "QUEM NÃO DEVE PARAR\n\n• T-score < -2.5\n• Fratura Prévia\n• Uso de Denosumabe"

# 4. A Mágica Final: Salvar
prs.save('Masterclass_Turner.pptx')

print("Mágica realizada. Arquivo criado.")

