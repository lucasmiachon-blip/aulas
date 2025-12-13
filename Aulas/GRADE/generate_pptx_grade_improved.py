import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# --- CONFIGURAÇÕES DE ESTILO (SÍNTESE & ELEGÂNCIA) ---
BG_COLOR = RGBColor(250, 248, 240)     # Bege Creme (Papel Antigo)
TITLE_COLOR = RGBColor(47, 79, 79)     # Dark Slate Gray (Sálvia Escuro)
TEXT_COLOR = RGBColor(80, 80, 80)      # Cinza Quente (Suave)
ACCENT_COLOR = RGBColor(160, 82, 45)   # Sienna (Terracota)

# Fontes (O PowerPoint usará se estiverem instaladas, senão usa padrão)
FONT_TITLE = 'Garamond'       # Clássica, acadêmica, elegante
FONT_BODY = 'Calibri Light'   # Limpa, moderna, legível


# --- MAPA DE IMAGENS ---
# Nomes exatos dos arquivos que você fez upload
IMAGENS = {
    "capa": "majestic_ancient_oak_tree_with_complex_exposed_roots_watercol_63974e8c-ac1a-49d4-869a-04d35982e582_0.png",
    "lente": "Watercolor_illustration_of_a_vintage_camera_lens_with_a_small_d419564b-7ca5-4eea-b963-56604a9efbdb_2.png",
    "bussolas": "Watercolor_painting_of_three_vintage_compasses_each_pointing__5cb9cfaf-26b0-4fe0-93e8-94fd0bc403a1_0.png",
    "nevoa": "Watercolor_painting_of_a_landscape_in_heavy_fog_indistinct_sh_6a6d886e-70bd-4780-a221-47c7265e01a2_0.png",
    "iceberg": "Watercolor_painting_of_an_iceberg_only_the_tip_is_visible_abo_d2ab112a-5d95-4e7c-96c9-119feb6b84fa_3.png",
    "tinta_diabetes": "op_of_golden_amber_ink_dissolving_in_clear_water_creating_sof_4e717e45-0d7c-4a7f-af74-02fd8ed13e12_2.png",
    "rios_lipides": "Aerial_watercolor_painting_of_two_soft_rivers_merging_into_on_1a3f996e-9601-47e7-bf47-b621b0966763_3.png",
    "folha_micro": "Macro_watercolor_painting_of_the_delicate_vein_structure_of_a_d39ecd70-5bda-4217-8349-f09e9556d327_0.png",
    "farol": "An_oil_painting_in_the_style_of_Claude_Monet._A_lighthouse_st_091c891a-43bd-4166-b801-be07761b5f58_0.png"
}

def aplicar_design_clean(slide):
    """Aplica o fundo bege minimalista"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR

def criar_slide_zen(prs, titulo, bullets, notas="", imagem_key=None):
    """
    Cria um slide focado em síntese:
    - Esquerda: Texto curto, poucas linhas.
    - Direita: Imagem grande e bela.
    """
    # Layout Blank para controle total
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    aplicar_design_clean(slide)

    # --- 1. IMAGEM (Lado Direito - Ocupa 50% da tela) ---
    if imagem_key and imagem_key in IMAGENS:
        img_path = IMAGENS[imagem_key]
        if os.path.exists(img_path):
            left = Inches(6.5)
            top = Inches(0) # Sangria total (topo ao chão) se possível, ou margem
            height = Inches(7.5)
            try:
                # Centralizar verticalmente ou ocupar lado direito
                pic = slide.shapes.add_picture(img_path, left, Inches(0.8), height=Inches(5.9))
                # Centralizar a imagem no espaço direito visualmente
                pic.left = Inches(7.0)
            except Exception as e:
                print(f"Erro imagem: {e}")

    # --- 2. TÍTULO (Elegante e Serifado) ---
    # Caixa de título à esquerda
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(6), Inches(1.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    p.text = titulo
    p.font.name = FONT_TITLE
    p.font.size = Pt(44) # Grande
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR

    # --- 3. BULLETS (Sintéticos) ---
    # Caixa de texto abaixo do título
    body_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.2), Inches(5.8), Inches(3.5))
    tf = body_box.text_frame
    tf.word_wrap = True
    
    # Adiciona cada bullet point
    if isinstance(bullets, list):
        for ponto in bullets:
            p = tf.add_paragraph()
            p.text = ponto
            p.font.name = FONT_BODY
            p.font.size = Pt(24) # Tamanho legível
            p.font.color.rgb = TEXT_COLOR
            p.space_after = Pt(20) # Respiro entre linhas
            # Simular bullet customizado (opcional, aqui usa o padrão)

    # --- 4. NOTAS DO ORADOR (A Profundidade) ---
    if notas:
        slide.notes_slide.notes_text_frame.text = notas


# --- INÍCIO DA GERAÇÃO ---
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

print("Gerando apresentação 'Soma Perfeita'...")

# --- CONTEÚDO (Curadoria de Sentidos) ---

# 1. Capa
criar_slide_zen(prs, 
    titulo="Dissecando a Evidência", 
    bullets=["O Sistema GRADE nas Diretrizes 2025", "Da incerteza à recomendação clínica."],
    notas="Bem-vindos. Hoje vamos navegar pelas raízes da decisão clínica.",
    imagem_key="capa")

# 2. Contexto
criar_slide_zen(prs, 
    titulo="A Tríade", 
    bullets=["SBC 2025 | ESC 2024 | AHA 2023", "Consenso clínico.", "Divergência metodológica."],
    notas="Três documentos massivos. Onde eles divergem, a chave não é a clínica, é o método GRADE.",
    imagem_key=None) # *Aqui você colará o print das capas*

# 3. O Tempo
criar_slide_zen(prs, 
    titulo="Prazo de Validade", 
    bullets=["A verdade científica perece.", "O GRADE é o filtro do tempo.", "Classe I ontem ➝ Classe III hoje."],
    notas="A ciência é viva. O que era verdade absoluta em 2010 pode ser erro médico hoje.",
    imagem_key=None)

# 4. A Mudança
criar_slide_zen(prs, 
    titulo="Pirâmide vs. GRADE", 
    bullets=["Adeus à hierarquia rígida.", "RCT ruim < Observacional forte.", "A qualidade define a posição."],
    notas="Esqueçam a pirâmide onde o RCT sempre ganha. No GRADE, a qualidade do desenho manda mais que o tipo de estudo.",
    imagem_key=None)

# 5. Vieses
criar_slide_zen(prs, 
    titulo="A Distorção da Lente", 
    bullets=["Risco de Viés (Bias).", "Sem cegamento, a verdade entorta.", "Falha de desenho = Falsa certeza."],
    notas="Assim como uma lente rachada deforma a paisagem, a falta de cegamento deforma o resultado.",
    imagem_key="lente")

# 6. Aplicação Viés
criar_slide_zen(prs, 
    titulo="O Peso do Desfecho", 
    bullets=["STICH (Cirurgia): Desfecho Duro (Morte).", "REVIVED (PCI): Desfecho Mole.", "Viés tolerável vs. Viés crítico."],
    notas="Por que a cirurgia mantém Classe I? Porque morte é difícil de mascarar. Já a PCI na IC sofreu com vieses em estudos antigos.",
    imagem_key=None) # *Colar gráfico STICH*

# 7. Inconsistência
criar_slide_zen(prs, 
    titulo="Nortes Diferentes", 
    bullets=["Inconsistência (Heterogeneidade).", "Quando os estudos não conversam.", "A certeza se dilui."],
    notas="Se uma bússola aponta pro Norte e outra pro Sul, você não confia em nenhuma. Isso acontece com os Betabloqueadores hoje.",
    imagem_key="bussolas")

# 8. Evidência Indireta
criar_slide_zen(prs, 
    titulo="O Alvo Errado", 
    bullets=["Evidência Indireta.", "Pergunta certa, população errada.", "Ex: ISCHEMIA excluiu Tronco."],
    notas="Não use o ISCHEMIA para justificar tratamento de Tronco. O tiro não acertou esse alvo.",
    imagem_key=None)

# 9. Imprecisão
criar_slide_zen(prs, 
    titulo="A Névoa Estatística", 
    bullets=["Imprecisão.", "Intervalos de Confiança amplos.", "Salva muito ou mata um pouco?"],
    notas="Quando o intervalo cruza a linha do nulo ou do dano, a névoa nos impede de ver a estrada.",
    imagem_key="nevoa")

# 10. Exemplo Imprecisão
criar_slide_zen(prs, 
    titulo="A Dúvida do Placebo", 
    bullets=["Omega-3 (Icosapent).", "REDUCE-IT vs STRENGTH.", "A dúvida gera recomendação fraca (IIb)."],
    notas="O placebo de óleo mineral atrapalhou a análise. A imprecisão derrubou a nota.",
    imagem_key=None)

# 11. Viés de Publicação
criar_slide_zen(prs, 
    titulo="A Ponta do Iceberg", 
    bullets=["Viés de Publicação.", "Resultados negativos somem.", "O milagre pode ser ilusão."],
    notas="Só vemos o que foi publicado. O que está submerso muitas vezes nega o tratamento.",
    imagem_key="iceberg")

# 12. O Bônus (Rate Up)
criar_slide_zen(prs, 
    titulo="Quando a Nota Sobe", 
    bullets=["Magnitude de Efeito.", "Gradiente Dose-Resposta.", "Observacional vira Certeza."],
    notas="Não precisamos de RCT para paraquedas. Às vezes o efeito é tão óbvio que a nota sobe.",
    imagem_key=None)

# 13. Diabetes (A Síntese Perfeita)
criar_slide_zen(prs, 
    titulo="A Viscosidade do Risco", 
    bullets=["Diabetes Mellitus.", "Magnitude: RR > 2x.", "Risco Equivalente à Doença Arterial."],
    notas="O GRADE reconhece o peso. O risco é tão denso que tratamos como prevenção secundária.",
    imagem_key="tinta_diabetes")

# 14. Aplicação Prática
criar_slide_zen(prs, 
    titulo="A Nova Sinfonia", 
    bullets=["Lípides: Estatina + Ezetimiba.", "A Confluência (Sinergia).", "Meta < 50 mg/dL com segurança."],
    notas="Rios que correm juntos chegam mais longe. A combinação supera a dose alta isolada.",
    imagem_key="rios_lipides")

# 15. Microcirculação
criar_slide_zen(prs, 
    titulo="O Detalhe Define", 
    bullets=["INOCA / Microcirculação.", "A anatomia não é tudo.", "Nervuras importam tanto quanto o tronco."],
    notas="A diretriz 2025 joga luz nas pequenas coisas. Diagnóstico funcional é vital.",
    imagem_key="folha_micro")

# 16. Encerramento
criar_slide_zen(prs, 
    titulo="A Visão Inteira", 
    bullets=["'Para ser grande, sê inteiro.'", "Nada teu exagera ou exclui.", "A ciência brilha alta."],
    notas="O poema de Pessoa resume o GRADE. Não exagerar o benefício, não excluir o risco. Obrigado.",
    imagem_key="farol")


# Salvar
file_name = 'Aula_GRADE_Sintese_Perfeita.pptx'
prs.save(file_name)
print(f"Sucesso! Apresentação salva em: {file_name}")